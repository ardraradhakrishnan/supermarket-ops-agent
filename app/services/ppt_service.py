import os
import io
import logging
from datetime import datetime
from sqlalchemy.orm import Session
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.services.base_service import BaseService
from app.services.analytics_service import AnalyticsService
from app.services.preference_service import PreferenceService

logger = logging.getLogger(__name__)

# Premium Color Palette
PRIMARY_COLOR = RGBColor(37, 99, 235)    # Blue
SECONDARY_COLOR = RGBColor(30, 41, 59)  # Dark slate
ACCENT_COLOR = RGBColor(248, 250, 252)  # Light gray
TEXT_DARK = RGBColor(15, 23, 42)       # Dark charcoal
TEXT_LIGHT = RGBColor(100, 116, 139)   # Slate gray
WHITE = RGBColor(255, 255, 255)

class PPTService(BaseService):
    """
    Handles PowerPoint presentation generation for supermarket business reporting.
    """

    def __init__(self, db: Session):
        super().__init__(db)
        self.analytics_service = AnalyticsService(db)
        self.pref_service = PreferenceService(db)

    def generate_sales_report_ppt(self, days: int = 30) -> str:
        """
        Generates a premium business report slide presentation.
        Returns the absolute filepath to the generated PPTX file.
        """
        store_name = self.pref_service.get_store_name()
        currency = self.pref_service.get_currency()
        
        # 1. Fetch data
        dashboard = self.analytics_service.get_dashboard_summary()
        top_products = self.analytics_service.get_top_selling_products(limit=5)
        low_stock = self.analytics_service.get_low_stock_report()
        inv_value = self.analytics_service.get_inventory_value()
        
        # Create output directory
        output_dir = os.path.abspath(os.path.join("data", "reports"))
        os.makedirs(output_dir, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = f"Report_{timestamp}.pptx"
        ppt_path = os.path.join(output_dir, ppt_filename)
        
        # Initialize Presentation
        prs = Presentation()
        # Set to 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        blank_slide_layout = prs.slide_layouts[6]
        
        # ==============================================================================
        # SLIDE 1: Title Slide (Dark Theme)
        # ==============================================================================
        slide1 = prs.slides.add_slide(blank_slide_layout)
        
        # Add background shape
        bg = slide1.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = SECONDARY_COLOR
        bg.line.fill.background()
        
        # Title text box
        title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{store_name} Operations Report"
        p.font.name = "Arial"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"30-Day Business Performance & Operations Insights"
        p2.font.name = "Arial"
        p2.font.size = Pt(22)
        p2.font.color.rgb = PRIMARY_COLOR
        p2.space_before = Pt(10)
        
        p3 = tf.add_paragraph()
        p3.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        p3.font.name = "Arial"
        p3.font.size = Pt(14)
        p3.font.color.rgb = TEXT_LIGHT
        p3.space_before = Pt(30)
        
        # ==============================================================================
        # SLIDE 2: Business Overview / KPI Dashboard
        # ==============================================================================
        slide2 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide2, "Business Overview Dashboard", "Key metrics and current operations status")
        
        # Grid of KPI Cards (4 Cards)
        kpis = [
            ("Total Products", f"{dashboard['total_products']}", "Catalog Count"),
            ("Active Customers", f"{dashboard['total_customers']}", "Registered Members"),
            ("Pending Khata Credit", f"{currency} {dashboard['pending_khata']:.2f}", "Outstanding Balance"),
            ("Est. Inventory Value", f"{currency} {inv_value['inventory_value']:.2f}", "Current Asset Value")
        ]
        
        card_width = Inches(2.6)
        card_height = Inches(3.0)
        card_top = Inches(2.2)
        spacing = Inches(0.4)
        left_margin = Inches(1.0)
        
        for i, (title, val, subtitle) in enumerate(kpis):
            left_pos = left_margin + i * (card_width + spacing)
            
            # Card Background
            card = slide2.shapes.add_shape(
                1, left_pos, card_top, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = ACCENT_COLOR
            card.line.color.rgb = PRIMARY_COLOR if "Khata" in title or "Value" in title else TEXT_LIGHT
            card.line.width = Pt(1.5)
            
            # Card Text
            tb = slide2.shapes.add_textbox(left_pos + Inches(0.1), card_top + Inches(0.2), card_width - Inches(0.2), card_height - Inches(0.4))
            ctf = tb.text_frame
            ctf.word_wrap = True
            
            cp1 = ctf.paragraphs[0]
            cp1.text = title.upper()
            cp1.font.name = "Arial"
            cp1.font.size = Pt(12)
            cp1.font.bold = True
            cp1.font.color.rgb = TEXT_LIGHT
            cp1.alignment = PP_ALIGN.CENTER
            
            cp2 = ctf.add_paragraph()
            cp2.text = val
            cp2.font.name = "Arial"
            cp2.font.size = Pt(26)
            cp2.font.bold = True
            cp2.font.color.rgb = PRIMARY_COLOR
            cp2.space_before = Pt(40)
            cp2.alignment = PP_ALIGN.CENTER
            
            cp3 = ctf.add_paragraph()
            cp3.text = subtitle
            cp3.font.name = "Arial"
            cp3.font.size = Pt(11)
            cp3.font.color.rgb = TEXT_LIGHT
            cp3.space_before = Pt(40)
            cp3.alignment = PP_ALIGN.CENTER

        # Add Low Stock Alert Banner at the bottom
        low_stock_count = dashboard["low_stock_products"]
        banner = slide2.shapes.add_shape(
            1, Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.8)
        )
        banner.fill.solid()
        if low_stock_count > 0:
            banner.fill.fore_color.rgb = RGBColor(254, 242, 242) # light red
            banner.line.color.rgb = RGBColor(239, 68, 68) # red
            alert_text = f"ATTENTION: There are {low_stock_count} products low on stock! Action required."
            text_color = RGBColor(153, 27, 27)
        else:
            banner.fill.fore_color.rgb = RGBColor(240, 253, 244) # light green
            banner.line.color.rgb = RGBColor(34, 197, 94) # green
            alert_text = "STATUS OK: All inventory levels are healthy. No immediate low stock warnings."
            text_color = RGBColor(22, 101, 52)
            
        btb = slide2.shapes.add_textbox(Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.6))
        btf = btb.text_frame
        bp = btf.paragraphs[0]
        bp.text = alert_text
        bp.font.name = "Arial"
        bp.font.size = Pt(14)
        bp.font.bold = True
        bp.font.color.rgb = text_color
        bp.alignment = PP_ALIGN.CENTER
        
        # ==============================================================================
        # SLIDE 3: Top Selling Products (Visual Chart Slide)
        # ==============================================================================
        slide3 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide3, "Top Selling Products", "Best-performing items in terms of units sold")
        
        if top_products:
            # Generate matplotlib chart
            products_names = [item["product"] for item in top_products]
            quantities = [float(item["quantity"] or 0) for item in top_products]
            
            # Reverse for horizontal bar chart
            products_names.reverse()
            quantities.reverse()
            
            fig, ax = plt.subplots(figsize=(6, 4))
            bars = ax.barh(products_names, quantities, color='#2563eb', edgecolor='#1e40af')
            
            # Styling chart
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#94a3b8')
            ax.spines['bottom'].set_color('#94a3b8')
            ax.tick_params(colors='#475569', labelsize=10)
            ax.set_title("Units Sold By Product", color='#1e293b', fontsize=12, fontweight='bold')
            
            # Label bars
            for bar in bars:
                width = bar.get_width()
                ax.text(width + 0.1, bar.get_y() + bar.get_height()/2, f'{int(width)}', 
                        va='center', ha='left', color='#0f172a', fontweight='bold', fontsize=9)
            
            plt.tight_layout()
            
            # Save chart to stream
            chart_stream = io.BytesIO()
            plt.savefig(chart_stream, format='png', dpi=200)
            chart_stream.seek(0)
            plt.close(fig)
            
            # Insert chart image into slide
            slide3.shapes.add_picture(chart_stream, Inches(1.0), Inches(2.2), width=Inches(5.5), height=Inches(4.2))
            
            # Add table on the right side
            rows, cols = len(top_products) + 1, 2
            left, top_tbl, width_tbl, height_tbl = Inches(7.0), Inches(2.2), Inches(5.333), Inches(3.8)
            table_shape = slide3.shapes.add_table(rows, cols, left, top_tbl, width_tbl, height_tbl)
            table = table_shape.table
            
            # Header cells
            table.cell(0, 0).text_frame.paragraphs[0].text = "Product Name"
            table.cell(0, 1).text_frame.paragraphs[0].text = "Quantity Sold"
            for col_idx in range(2):
                cell = table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = SECONDARY_COLOR
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                
            # Table values
            # Reverse back to show highest first
            top_products.reverse()
            for idx, item in enumerate(top_products):
                table.cell(idx+1, 0).text_frame.paragraphs[0].text = str(item["product"])
                table.cell(idx+1, 1).text_frame.paragraphs[0].text = f"{int(item['quantity'] or 0)} units"
                
                # Format cell text
                for c_idx in range(2):
                    cell = table.cell(idx+1, c_idx)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)
                    p.font.color.rgb = TEXT_DARK
                    
        else:
            # Fallback text box if no sales yet
            no_data_box = slide3.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(2.0))
            tf = no_data_box.text_frame
            p = tf.paragraphs[0]
            p.text = "No sales transactions recorded in the system yet. Seed data or record sales to populate chart."
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            
        # ==============================================================================
        # SLIDE 4: Operational Tasks & Low Stock Report
        # ==============================================================================
        slide4 = prs.slides.add_slide(blank_slide_layout)
        self._add_slide_header(slide4, "Inventory Health & Alert List", "Details of products requiring replenishment")
        
        # Check low stock items
        if low_stock:
            # Create table
            rows_ls = min(len(low_stock) + 1, 7) # Limit to 6 items to fit screen
            left_ls, top_ls, width_ls, height_ls = Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0)
            ls_table_shape = slide4.shapes.add_table(rows_ls, 3, left_ls, top_ls, width_ls, height_ls)
            ls_table = ls_table_shape.table
            
            # Setup headers
            headers = ["Product Name", "Current Stock", "Status"]
            for col_idx, h_text in enumerate(headers):
                cell = ls_table.cell(0, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY_COLOR
                p = cell.text_frame.paragraphs[0]
                p.text = h_text
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = WHITE
                
            for idx in range(rows_ls - 1):
                item = low_stock[idx]
                ls_table.cell(idx+1, 0).text_frame.paragraphs[0].text = str(item["product"])
                ls_table.cell(idx+1, 1).text_frame.paragraphs[0].text = f"{int(item['stock'])} units"
                ls_table.cell(idx+1, 2).text_frame.paragraphs[0].text = "CRITICAL: REORDER NOW"
                
                # Format cell text
                for c_idx in range(3):
                    cell = ls_table.cell(idx+1, c_idx)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(12)
                    p.font.color.rgb = TEXT_DARK
                    if c_idx == 2:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(185, 28, 28) # Red alert text
        else:
            # No low stock alert
            ok_box = slide4.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(2.0))
            tf = ok_box.text_frame
            p = tf.paragraphs[0]
            p.text = "All products are well above their reorder thresholds!\nNo inventory action items."
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(22, 101, 52)
            p.alignment = PP_ALIGN.CENTER
            
        # Save presentation
        try:
            prs.save(ppt_path)
            logger.info(f"PowerPoint report successfully generated: {ppt_path}")
            return ppt_path
        except Exception as e:
            logger.error(f"Failed to save PPTX file: {e}")
            raise e

    def _add_slide_header(self, slide, title: str, subtitle: str):
        """
        Helper method to add a premium uniform header to slides.
        """
        # Header title
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Arial"
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(4)
        
        # Slide dividing line
        line = slide.shapes.add_shape(
            1, Inches(1.0), Inches(1.7), Inches(11.333), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = PRIMARY_COLOR
        line.line.fill.background()
